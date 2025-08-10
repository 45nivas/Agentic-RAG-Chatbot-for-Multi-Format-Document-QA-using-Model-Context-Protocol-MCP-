import os
from typing import List

def chunk_text(text: str, chunk_size: int = 500) -> List[str]:
    if not text or not text.strip():
        return []
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size) if text[i:i+chunk_size].strip()]

def parse_pdf(path: str) -> List[str]:
    try:
        from PyPDF2 import PdfReader
        reader = PdfReader(path)
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
        return chunk_text(text)
    except Exception as e:
        return [f"Error parsing PDF: {str(e)}"]

def parse_docx(path: str) -> List[str]:
    try:
        from docx import Document
        doc = Document(path)
        text = "\n".join(p.text for p in doc.paragraphs)
        return chunk_text(text)
    except Exception as e:
        return [f"Error parsing DOCX: {str(e)}"]

def parse_pptx(path: str) -> List[str]:
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
            slides.append(f"slide {i+1}: {slide_text}")
        return chunk_text("\n".join(slides))
    except Exception as e:
        return [f"Error parsing PPTX: {str(e)}"]

def parse_csv(path: str) -> List[str]:
    try:
        import pandas as pd
        df = pd.read_csv(path)
        return chunk_text(df.to_csv(index=False))
    except Exception as e:
        return [f"Error parsing CSV: {str(e)}"]

def parse_txt_md(path: str) -> List[str]:
    try:
        with open(path, encoding="utf-8") as f:
            text = f.read()
        return chunk_text(text)
    except Exception as e:
        return [f"Error parsing TXT/MD: {str(e)}"]

def parse_document(path: str) -> List[str]:
    try:
        if not os.path.exists(path):
            return [f"File not found: {path}"]
        
        ext = os.path.splitext(path)[-1].lower()
        if ext == ".pdf":
            return parse_pdf(path)
        elif ext == ".docx":
            return parse_docx(path)
        elif ext == ".pptx":
            return parse_pptx(path)
        elif ext == ".csv":
            return parse_csv(path)
        elif ext in [".txt", ".md"]:
            return parse_txt_md(path)
        else:
            return [f"Unsupported file type: {ext}"]
    except Exception as e:
        return [f"Error parsing document {path}: {str(e)}"]
