# -*- coding: utf-8 -*-
import sys
import io
# Set UTF-8 encoding for Windows console
if sys.platform == 'win32':
    sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')
    sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8')

from flask import Flask, render_template, request, jsonify, session
import os
import logging
from datetime import datetime
from werkzeug.utils import secure_filename
from dotenv import load_dotenv

# Add parent directory to path so we can import agents/
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..'))

# Multi-Agent Architecture with MCP Protocol
from agents.coordinator_agent import CoordinatorAgent
from agents.ingestion_agent import IngestionAgent
from agents.retrieval_agent import RetrievalAgent
from agents.llm_response_agent import LLMResponseAgent
from agents.mcp import MCPMessage
from agents.embedding_utils import EmbeddingStore

# Advanced AI Stack
from sentence_transformers import SentenceTransformer
import chromadb
from chromadb.config import Settings
import numpy as np
print("✅ Sentence Transformers + ChromaDB - Advanced AI Stack!")

import google.generativeai as genai
import json

load_dotenv(os.path.join(os.path.dirname(__file__), '..', '.env'))

app = Flask(__name__)
app.secret_key = os.environ.get('SECRET_KEY', 'dev-secret-advanced')
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024
app.config['DEBUG'] = True

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure Gemini
GEMINI_API_KEY = os.environ.get('GEMINI_API_KEY')
model = None
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    model = genai.GenerativeModel('models/gemini-2.5-flash')
    logger.info("✅ Gemini AI configured")
else:
    logger.warning("⚠️ No GEMINI_API_KEY found - responses will be limited")


# ============================================================
# Multi-Agent RAG System using Model Context Protocol (MCP)
# ============================================================
# Agents:
#   1. IngestionAgent  - Parses documents into chunks (MCPMessage)
#   2. RetrievalAgent  - Embeds chunks + semantic search (MCPMessage)
#   3. LLMResponseAgent - Generates answers via Gemini (MCPMessage)
#   4. CoordinatorAgent - Orchestrates the full pipeline (MCPMessage)
# All inter-agent communication uses MCPMessage protocol
# ============================================================

class AgenticRAG:
    """Multi-Agent RAG with Sentence Transformers + ChromaDB + MCP Protocol"""

    def __init__(self):
        # Initialize the 4-agent system
        self.coordinator = CoordinatorAgent()
        self.ingestion_agent = IngestionAgent()
        self.retrieval_agent = RetrievalAgent()
        self.llm_agent = LLMResponseAgent()

        # Initialize embedding store (shared by retrieval agent)
        self.embedding_store = EmbeddingStore(persist_path="./chroma_advanced")
        self.retrieval_agent.vector_store = self.embedding_store

        self.file_paths = []  # Track uploaded file paths
        self.mcp_trace = []   # Track all MCP messages for debugging

        logger.info("✅ 4-Agent System initialized with MCP protocol")
        logger.info("  → IngestionAgent (document parsing)")
        logger.info("  → RetrievalAgent (Sentence Transformers 384D + ChromaDB + HNSW)")
        logger.info("  → LLMResponseAgent (Gemini AI)")
        logger.info("  → CoordinatorAgent (orchestrator)")
        logger.info("🚀 Advanced Agentic RAG initialized!")

    def ingest_document(self, file_path, filename):
        """Use IngestionAgent to parse document → returns MCPMessage"""
        # IngestionAgent parses the file and returns an MCPMessage
        ingest_msg = self.ingestion_agent.parse_documents([file_path])
        self.mcp_trace.append(ingest_msg.to_dict())

        chunks = ingest_msg.payload.get("chunks", [])
        logger.info(f"📨 MCP [{ingest_msg.sender} → {ingest_msg.receiver}] "
                     f"type={ingest_msg.type} | {len(chunks)} chunks from {filename}")

        if not chunks:
            logger.warning(f"⚠️ No chunks extracted from {filename}")
            return False

        # Add chunks to RetrievalAgent's embedding store
        self.embedding_store.add_chunks(chunks)
        logger.info(f"✅ {len(chunks)} chunks embedded via Sentence Transformers (384D) for {filename}")

        return True

    def query(self, question):
        """Full agent pipeline: Retrieve → LLM Response via MCP messages"""

        # Step 1: RetrievalAgent searches for relevant chunks
        retrieval_msg = self.retrieval_agent.embed_and_retrieve([], question)
        self.mcp_trace.append(retrieval_msg.to_dict())

        context = retrieval_msg.payload.get("retrieved_context", [])
        max_similarity = retrieval_msg.payload.get("max_similarity", 0.0)
        threshold_met = retrieval_msg.payload.get("threshold_met", False)

        logger.info(f"📨 MCP [{retrieval_msg.sender} → {retrieval_msg.receiver}] "
                     f"type={retrieval_msg.type} | {len(context)} results, "
                     f"max_sim={max_similarity:.3f}")

        if not context:
            return {
                'answer': "No relevant information found in the uploaded documents.",
                'similarity': 0.0,
                'method': 'RetrievalAgent → No results',
                'mcp_messages': len(self.mcp_trace)
            }

        # Step 2: LLMResponseAgent generates answer using retrieved context
        llm_msg = self.llm_agent.generate_response(
            context, question, threshold_met, max_similarity
        )
        self.mcp_trace.append(llm_msg.to_dict())

        answer = llm_msg.payload.get("answer", "Could not generate response.")

        logger.info(f"📨 MCP [{llm_msg.sender} → {llm_msg.receiver}] "
                     f"type={llm_msg.type} | response generated")

        return {
            'answer': answer,
            'similarity': max_similarity,
            'method': 'IngestionAgent → RetrievalAgent → LLMResponseAgent (MCP)',
            'mcp_messages': len(self.mcp_trace),
            'threshold_met': threshold_met
        }

    def clear(self):
        """Clear all data and reset agents"""
        self.embedding_store.clear()
        self.retrieval_agent.vector_store = self.embedding_store
        self.file_paths = []
        self.mcp_trace = []
        logger.info("🗑️ All agent data cleared")

    def get_mcp_trace(self):
        """Return all MCP messages for debugging/display"""
        return self.mcp_trace


# Initialize the Multi-Agent RAG system
rag = AgenticRAG()


# File processing helpers
def process_file(file_path, filename):
    """Extract text from uploaded file"""
    try:
        ext = filename.rsplit('.', 1)[-1].lower()

        if ext == 'pdf':
            import PyPDF2
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                return " ".join([page.extract_text() for page in reader.pages])

        elif ext == 'docx':
            import docx
            doc = docx.Document(file_path)
            return " ".join([para.text for para in doc.paragraphs])

        elif ext == 'pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            slides_text = []
            for i, slide in enumerate(prs.slides):
                slide_text = " ".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                slides_text.append(f"Slide {i+1}: {slide_text}")
            return "\n".join(slides_text)

        elif ext == 'csv':
            import pandas as pd
            df = pd.read_csv(file_path)
            return df.to_string()

        elif ext in ['txt', 'md']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                return f.read()

        return ""
    except Exception as e:
        logger.error(f"Processing error: {e}")
        return ""


# Flask routes
UPLOAD_FOLDER = 'uploads'
ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'csv', 'txt', 'md'}
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER


@app.route('/')
def index():
    return render_template('index.html')


@app.route('/upload', methods=['POST'])
def upload():
    """Upload documents → IngestionAgent parses via MCP"""
    try:
        # Clear previous data before new upload
        logger.info("🗑️ Clearing previous data before new upload...")
        rag.clear()
        session.clear()
        logger.info("✅ Previous data cleared")

        files = request.files.getlist('files')
        processed = []
        mcp_messages_generated = 0

        for file in files:
            if file and file.filename:
                filename = secure_filename(file.filename)
                if filename.rsplit('.', 1)[-1].lower() in ALLOWED_EXTENSIONS:
                    # Save temporarily
                    file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                    file.save(file_path)

                    # Use IngestionAgent via MCP to parse document
                    if rag.ingest_document(file_path, filename):
                        processed.append(filename)
                        mcp_messages_generated = len(rag.get_mcp_trace())

                    # Cleanup temp file
                    try:
                        os.remove(file_path)
                    except:
                        pass

        if processed:
            return jsonify({
                'success': True,
                'message': f'Processed {len(processed)} files via IngestionAgent (MCP)',
                'files': processed,
                'method': 'IngestionAgent → EmbeddingStore (Sentence Transformers 384D + ChromaDB)',
                'mcp_messages': mcp_messages_generated,
                'cleared': True
            })
        else:
            return jsonify({'error': 'No files processed'}), 400

    except Exception as e:
        logger.error(f"❌ Upload error: {e}")
        return jsonify({
            'success': False,
            'error': str(e),
            'method': 'Error in IngestionAgent'
        }), 500


@app.route('/chat', methods=['POST'])
def chat():
    """Chat → RetrievalAgent + LLMResponseAgent via MCP pipeline"""
    logger.info("📨 Chat request received")
    try:
        data = request.get_json()
        question = data.get('message', '').strip()
        logger.info(f"❓ Question: {question}")

        if not question:
            return jsonify({'error': 'No question provided'}), 400

        # Run the full agent pipeline via MCP
        result = rag.query(question)

        # Update conversation history
        if 'conversation' not in session:
            session['conversation'] = []

        session['conversation'].append({
            'question': question,
            'answer': result['answer'],
            'timestamp': datetime.now().isoformat(),
            'similarity': result['similarity'],
            'method': result['method']
        })

        return jsonify({
            'success': True,
            'response': result['answer'],
            'method': result['method'],
            'similarity': round(result['similarity'], 3),
            'mcp_messages': result['mcp_messages']
        })

    except Exception as e:
        logger.error(f"❌ Chat error: {e}")
        return jsonify({'error': 'Chat processing failed'}), 500


@app.route('/clear', methods=['POST'])
def clear():
    """Clear all agent data and conversation"""
    try:
        session.clear()
        rag.clear()

        return jsonify({
            'success': True,
            'message': 'All agent data and conversation cleared',
            'details': 'ChromaDB collection reset, MCP trace cleared'
        })

    except Exception as e:
        logger.error(f"❌ Clear error: {e}")
        return jsonify({
            'success': False,
            'error': f'Failed to clear: {str(e)}'
        }), 500


@app.route('/health', methods=['GET'])
def health():
    """Health check with agent status"""
    return jsonify({
        'status': 'healthy',
        'agents': {
            'ingestion': 'active',
            'retrieval': 'active',
            'llm_response': 'active',
            'coordinator': 'active'
        },
        'embedding': 'Sentence Transformers (384D)',
        'vector_db': 'ChromaDB + HNSW',
        'llm': 'Gemini AI',
        'protocol': 'MCP (Model Context Protocol)',
        'mcp_messages_total': len(rag.get_mcp_trace())
    })


if __name__ == '__main__':
    logger.info("🚀 Agentic RAG starting with MCP Protocol!")
    logger.info("   Sentence Transformers + ChromaDB + HNSW + Gemini AI")
    logger.info("   4 Agents: Ingestion → Retrieval → LLM → Coordinator")
    app.run(debug=True, host='0.0.0.0', port=5000)
